
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])

title = slide.shapes.title
title.text = "Meeting Summary - Financial Results and Strategies"

subtitle = slide.placeholders[1]
subtitle.text = "Slide 1"

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Agenda"

tf = body_shape.text_frame
tf.text = "Introduction"

p = tf.add_paragraph()
p.text = "Financial Results"
p.level = 1

p = tf.add_paragraph()
p.text = "Strategies for the Upcoming Year"
p.level = 1

p = tf.add_paragraph()
p.text = "Q&A"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Introduction"

tf = body_shape.text_frame
tf.text = "Welcome and Introduction by Giacomo"

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Financial Results"

tf = body_shape.text_frame
tf.text = "Presentation by Andrea"

p = tf.add_paragraph()
p.text = "Net Profit: 10 million euros, 5% increase YoY"
p.level = 1

p = tf.add_paragraph()
p.text = "Improvements in net interest margin, efficiency ratio, and return on equity"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Key Factors Contributing to Financial Results"

tf = body_shape.text_frame
tf.text = "Attracting new customers in the small and medium-sized businesses segment"

p = tf.add_paragraph()
p.text = "Careful cost and risk management"
p.level = 1

p = tf.add_paragraph()
p.text = "Investment in digital innovation"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Strategies for the Upcoming Year"

tf = body_shape.text_frame
tf.text = "Consolidating market position"

p = tf.add_paragraph()
p.text = "Increasing profitability"
p.level = 1

p = tf.add_paragraph()
p.text = "Enhancing reputation"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Priorities for the Upcoming Year"

tf = body_shape.text_frame
tf.text = "Continuing focus on customer needs"

p = tf.add_paragraph()
p.text = "Developing new products and services"
p.level = 1

p = tf.add_paragraph()
p.text = "Strengthening online presence"
p.level = 1

p = tf.add_paragraph()
p.text = "Supporting ecological and social transition"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Challenges"

tf = body_shape.text_frame
tf.text = "Increasing competition"

p = tf.add_paragraph()
p.text = "Adapting to new regulations"
p.level = 1

p = tf.add_paragraph()
p.text = "Managing economic uncertainty related to the pandemic"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Financing of Digital Innovation Projects"

tf = body_shape.text_frame
tf.text = "Question by Director E"

p = tf.add_paragraph()
p.text = "Response by Andrea:"
p.level = 1

p = tf.add_paragraph()
p.text = "Current financial position is strong"
p.level = 2

p = tf.add_paragraph()
p.text = "Good access to capital markets"
p.level = 2

p = tf.add_paragraph()
p.text = "No plans to increase capital or seek external loans"
p.level = 2

p = tf.add_paragraph()
p.text = "Continuous monitoring of financial situation and evaluation of opportunities"
p.level = 2

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Conclusion"

tf = body_shape.text_frame
tf.text = "Summary of key points"

p = tf.add_paragraph()
p.text = "Thank you and good day by Giacomo"
p.level = 1

p = tf.add_paragraph()
p.text = "Thank you and extension of thanks by Andrea"
p.level = 1

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Q&A"

tf = body_shape.text_frame
tf.text = "Opportunity for participants to ask additional questions or seek clarification"

slide = prs.slides.add_slide(prs.slide_layouts[1])

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Closing"

tf = body_shape.text_frame
tf.text = "Final thank you and closing remarks by Giacomo"

prs.save('presentation.pptx')


