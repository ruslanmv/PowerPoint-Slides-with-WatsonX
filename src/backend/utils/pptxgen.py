from utils.generator import Generator
from ibm_watson_machine_learning.foundation_models.utils.enums import ModelTypes
import requests
import logging

class pptxgenerator(Generator):
    code_prompt = f"""You are really good at writing python code and using the library python-pptx.
Write the code in python to create a single slide for powerpoint presentation given the text.
Do not add notes and skip comments on the results. 
All the images added have the name image.png.
Do not add comments of the answer.
Do not explain the code.
Answer only with  Python code.
            """          
    def __init__(self):
        super().__init__(ModelTypes.LLAMA_2_70B_CHAT)
    def generate_code(self, text,n):
        with open('documentation.txt', 'r') as file:
            documentation = file.read()


        prefix_code=f"""from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
prs = Presentation()    
#Slide {n} :
slide = prs.slides.add_slide("""        
        suffix_code=f"""prs.save('slide_{n}.pptx')"""        
        logging.info("Sending prompt for code")
        inst_prompt = f"""<s>[INST] <<SYS>>  
        {pptxgenerator.code_prompt}
        You are doing the slide number {n}.
        You can start the code with:{prefix_code}.
        and end the code with:{suffix_code}.
        Give a complete python code ready to run.
        <</SYS>> 
        Create the python code given this text: {text}
        Documentation:{documentation}
        [/INST]
        """
        print("Prompt:",inst_prompt)
        result = self.llm_model.generate(
            inst_prompt)['results'][0]['generated_text']
        logging.info("Obtained result")
        return result
    

    
def extract_slides(text):
    slides = []
    lines = text.split("\n")
    current_slide = ""

    for line in lines:
        if line.startswith("Slide"):
            if current_slide:
                slides.append(current_slide)
                current_slide = ""
            current_slide = line
        elif current_slide:
            current_slide += "\n" + line

    if current_slide:
        slides.append(current_slide)
    
    return slides